import copy

from src import common_module as cm
from src.common.constants import SystemConstants, DbTypeConstants
from src.common.utils import DateUtils, SystemUtils, SqlUtils
from src.analysis_extend_target import OracleTarget

import pandas as pd

from resources.config_manager import Config

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pptx.table import Table, _Row, _Cell
from copy import deepcopy
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData


class ReportPpt(cm.CommonModule):

    def __init__(self, logger):
        super().__init__(logger=logger)
        self.ot: OracleTarget = None
        self.presentiton_path = 'report10.pptx'
        self.presentation = Presentation(self.presentiton_path)

    def main_process(self):

        self.logger.info("Report PPT")
        self._insert_extend_target_data()
        self._check_filename()

    def _insert_extend_target_data(self):
        """
        DB 확장 타겟(실제 분석 대상 DB) 데이터 저장 함수
        """
        extend_target_repo_list = self.config["maxgauge_repo"].get("extend_target_repo", [])

        for extend_target_repo in extend_target_repo_list:
            extend_target_repo["analysis_target_type"] = self.config["maxgauge_repo"]["analysis_target_type"]  # oracle

            if str(self.config["maxgauge_repo"]["analysis_target_type"]).lower() == DbTypeConstants.ORACLE:
                if self.ot is None:
                    self.ot = OracleTarget(self.logger, self.config)

                self.ot.set_extend_target_config(extend_target_repo)
                self.ot.init_process()

    def _extract_metric_name_list(self):

        metric_name_list = [
            "Host CPU Utilization (%)",
            "Average Active Sessions",
            "Executions Per Sec",
            "User Transaction Per Sec",
            "Logical Reads Per Sec",
            "Physical Reads Per Sec",
            "Hard Parse Count Per Sec",
        ]

        config_report = Config("report").get_config()
        metric_name_list.extend(config_report["sys_metric"])

        return metric_name_list


    def _convert_sql_to_df(self, sql_path, filename):
        """
        sql query문을 dataframe 형태로 변환
        """
        metric_name_list = self._extract_metric_name_list()
        unpack_metric_name_list = str(metric_name_list)[1:-1]

        s_date, e_date = DateUtils.get_each_date_by_interval2(
            self.config["args"]["s_date"], self.config["args"]["interval"], arg_fmt="%Y-%m-%d"
        )
        date_dict = {"StartDate": s_date, "EndDate": e_date, "Metric_Name": unpack_metric_name_list}

        query = SystemUtils.get_file_content_in_path(sql_path, filename + ".txt")
        date_query = SqlUtils.sql_replace_to_dict(query, date_dict)

        for df in self.ot.get_data_by_query(date_query):
            df.columns = [i.upper() for i in df.columns]
            return df

    def _extract_instance_num_df(self, df, category_name):
        """
        instance_num dataframe 추출
        """
        instance_num_df_list =[]
        for instance_num in df[category_name].unique():
            instance_num_df = df[df[category_name] == instance_num]
            instance_num_df_list.append(instance_num_df)
        return instance_num_df_list

    def _check_filename(self):

        sql_path = f"{self.config['home']}/" + SystemConstants.CHART_SQL
        txt_file_list = SystemUtils.get_filenames_from_path(sql_path)

        for file in txt_file_list:
            filename = file.split(".")[0]

            if filename == "TOP_N_Wait_Events":
                df = self._convert_sql_to_df(sql_path, filename)
                self._extract_table_data_in_ppt(df, 'TOP-N Wait Events', 'INSTANCE_NUMBER')

            if filename == "TOP_1_Wait_Events":
                df = self._convert_sql_to_df(sql_path, filename)
                self._extract_table_data_in_ppt2(df, 'TOP Wait Events – log file sync', 'INSTANCE_NUMBER')

            # if filename == "TOP_2_Wait_Events":
            #     df = self._convert_sql_to_df(sql_path, filename)
            #     self._extract_table_data_in_ppt2(df, 'TOP Wait Events – log file sync')

            # if filename == "TOP_1_Schema_SQL":
            #     df = self._convert_sql_to_df(sql_path,filename)
            #     self._extract_table_data_in_ppt2(df, 'TOP Schema & SQL')

            if filename == "Literal_SQL":
                df = self._convert_sql_to_df(sql_path,filename)
                self._extract_table_data(df, '성능 분석 – Literal SQL 점검')

    def _extract_table_data(self, df, text_frame_text):

        df_list=[]
        df_list.append(df)
        print('df_list',df_list)

        for idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text_frame.text == text_frame_text:
                    ppt_df_list = self._extract_ppt_df(slide)
                    result_list = self._compare_data(df_list, ppt_df_list)
                    print('result_list',result_list)
                    shape_list = self._extract_shape_list(slide, mso_type=MSO_SHAPE_TYPE.TABLE)
                    self._insert_data_into_ppt_table(shape_list, result_list)

    def _extract_table_data_in_ppt2(self, df, text_frame_text, cn):

        instance_df_list = self._extract_instance_num_df(df, category_name=cn)

        instance_df_columns = list(df)
        instance_df_columns.remove('INSTANCE_NUMBER')

        total_ppt_df_list=[]
        for idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text_frame.text == text_frame_text:
                    total_ppt_df_list.append(self._extract_ppt_df2(slide))

        shape_list=[]
        ppt_df_list=[]
        for i in total_ppt_df_list:
            for c in i:
                if instance_df_columns == list(c[1]):
                    shape_list.append(c[0])
                    ppt_df_list.append(c[1])

        result_list = self._compare_data(instance_df_list, ppt_df_list)
        self._insert_data_into_ppt_table(shape_list, result_list)

    def _extract_ppt_df2(self,slide):

        ppt_df_list=[]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                row_count = len(shape.table.rows)
                col_count = len(shape.table.columns)
                for _r in range(0, row_count):
                    row = []
                    for _c in range(0, col_count):
                        cell = shape.table.cell(_r, _c)
                        row.append(cell.text)
                    table_data.append(row)

                df_temp = pd.DataFrame(columns=table_data[0], data=table_data[1:])
                result = [shape, df_temp]
                ppt_df_list.append(result)

        return ppt_df_list

    def _extract_table_data_in_ppt(self,df,text_frame_text, cn):
        """

        """
        instance_df_list = self._extract_instance_num_df(df, category_name=cn)

        for idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text_frame.text == text_frame_text:
                    ppt_df_list = self._extract_ppt_df(slide)
                    result_list = self._compare_data(instance_df_list, ppt_df_list)
                    shape_list = self._extract_shape_list(slide, mso_type=MSO_SHAPE_TYPE.TABLE)
                    self._insert_data_into_ppt_table(shape_list, result_list)

    def _check_instance_num_and_slide_num(self,df,text_frame_text):
        instance_df_list = self._extract_instance_num_df(df)
        slide_list = []

        for idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and shape.text_frame.text == text_frame_text:
                    result_tuple = (idx, slide)
                    slide_list.append(result_tuple)

        if len(instance_df_list) == len(slide_list):
            print("slide 같음")
            self._insert_data_into_ppt2(df,slide_list)

        if len(instance_df_list) > len(slide_list):
            print("slide 추가")

            for idx, slide in slide_list:
                source_slide = self.presentation.slides[idx] # 복사하려는 슬라이드 정의
                slide_layout = source_slide.slide_layout
                copied_slide = self.presentation.slides.add_slide(slide_layout)

                title = copied_slide.placeholders[0]
                sp_title = title.element
                sp_title.getparent().remove(sp_title)

                subtitle = copied_slide.placeholders[1]
                sp = subtitle.element
                sp.getparent().remove(sp)

                for shape in source_slide.shapes:
                    if shape.has_chart:  # Check if the shape is a chart/graph

                        original_chart = shape.chart
                        category_labels = [c.label for c in original_chart.plots[0].categories]

                        chart_data = CategoryChartData()

                        for i, series in enumerate(original_chart.plots[0].series):
                            print(series.values)
                            chart_data.categories = category_labels
                            chart_data.add_series(series.name,series.values)

                        new_chart = copied_slide.shapes.add_chart(
                            XL_CHART_TYPE.LINE,
                            shape.left,
                            shape.top,
                            shape.width,
                            shape.height,
                            chart_data
                        )

                        new_chart.chart_style = original_chart.chart_style

                        # print("chart???", shape.chart)
                        # print(shape.shape_id)
                        # new_chart = copy.deepcopy(shape.chart)
                        # copied_slide.shapes._spTree.insert_element_before(new_chart.element, 'p:extLst')
                        # Handle chart copying manually

                    else:
                        # 차트를 제외한 데이터 copy
                        el = shape.element
                        newel = copy.deepcopy(el)
                        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

                self.presentation.slides._sldIdLst.insert(idx + 1, self.presentation.slides._sldIdLst[-1])
                del self.presentation.slides._sldIdLst[-1]

            self.presentation.save('plfi2.pptx')


        if len(instance_df_list) < len(slide_list):
            print(f"instance_df_list 개수: {len(instance_df_list)}, slide_list 개수 :{len(slide_list)}")
            remove_count = len(slide_list)-len(instance_df_list)
            xml_slides = self.presentation.slides._sldIdLst
            first_number = slide_list[0][0]
            last_number = slide_list[-1][0]
            slides = list(xml_slides)[first_number:last_number+1] #19, 20, 21번째 slide 출력

            remove_indices = [idx for idx in range(len(slide_list) - 1, len(slide_list) - 1 - remove_count, -1)]

            for idx in remove_indices:
                print(slides[idx])
                xml_slides.remove(slides[idx])
            self.presentation.save(f'delete_slide_{remove_count}.pptx')

    def _insert_data_into_ppt_table(self, shape_list, result_list):
        """

        """
        print(shape_list, result_list)
        for shape, result in zip(shape_list, result_list):
            row_count = len(shape.table.rows)

            for row_index, row_data in enumerate(result.values):
                if row_index + 1 < row_count:
                    for col_index, cell_data in enumerate(row_data):
                        cell = shape.table.cell(row_index + 1, col_index)
                        cell.text_frame.text = str(cell_data)
                        text_frame = cell.text_frame
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(8)

                if row_index + 1 >= row_count:
                    print("eeeeeeeeee", row_index, row_data, shape.table)
                    self.add_row(shape.table, row_data)

            self.presentation.save('231023_5.pptx')

    def add_row(self, table:Table, row_data) -> _Row:
        """

        """
        new_row = deepcopy(table._tbl.tr_lst[-1])

        for col_index, tc in enumerate(new_row.tc_lst):
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = str(row_data[col_index])  # Set cell text from row_data
            text_frame = cell.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(8)

        table._tbl.append(new_row)
        # return table._tbl
        # return table.rows[-1]
    def _extract_ppt_df(self,slide):

        ppt_df_list=[]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_data = []
                row_count = len(shape.table.rows)
                col_count = len(shape.table.columns)
                for _r in range(0, row_count):
                    row = []
                    for _c in range(0, col_count):
                        cell = shape.table.cell(_r, _c)
                        row.append(cell.text)
                    table_data.append(row)

                df_temp = pd.DataFrame(columns=table_data[0], data=table_data[1:])
                ppt_df_list.append(df_temp)

        return ppt_df_list

    def _extract_shape_list(self, slide, mso_type):
        shape_list = []
        for shape in slide.shapes:
            if shape.shape_type == mso_type:
                shape_list.append(shape)
        return shape_list

    def _compare_data(self,instance_df_list, ppt_df_list):
        """

        """

        zip_list = list(zip(instance_df_list, ppt_df_list))
        result_list = []

        for zip_df in zip_list:
            instance_num_df = zip_df[0]
            ppt_df = zip_df[1]
            result = instance_num_df[instance_num_df.columns.intersection(ppt_df.columns)]
            result_list.append(result)

        return result_list



