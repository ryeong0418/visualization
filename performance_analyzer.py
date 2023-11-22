import copy

from src import common_module as cm
from src.common.constants import SystemConstants, DbTypeConstants
from src.common.utils import DateUtils, SystemUtils, SqlUtils
from src.analysis_extend_target import OracleTarget
from src.ppt.ppt_writer import SlideManager

import pandas as pd
import re

from resources.config_manager import Config

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt
from pptx.enum.chart import  XL_DATA_LABEL_POSITION
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.chart.data import CategoryChartData
from pptx.table import _Cell
from pptx.util import Inches
from sql.common_sql import PptSql


class PerformanceAnalyzer(cm.CommonModule):

    FONT_NAME = "나눔스퀘어 네오 Light"

    SMALL_FONT_SIZE = Pt(8)
    DEFAULT_FONT_SIZE = Pt(10)
    MEDIUM_FONT_SIZE = Pt(11)
    LARGE_FONT_SIZE = Pt(24)

    BASE_LINE_SPACE = 1.0
    LONG_LINE_SPACE = 1.5

    BOLD_TRUE = True
    BOLD_FALSE = False

    YELLOW_COLOR = RGBColor(255, 217, 102)
    WHILE_COLOR = RGBColor(255, 255, 255)
    GREEN_COLOR = RGBColor(112, 173, 71)
    LIGHT_RED_COLOR = RGBColor(198, 29, 81)
    BLUE_GRAY_COLOR = RGBColor(65, 84, 110)
    DARK_GRAY_COLOR = RGBColor(89, 89, 89)
    LIGHT_BLUE_COLOR = RGBColor(61, 129, 246)
    BLACK_COLOR = RGBColor(0,0,0)

    TIME_MODEL = "#TIME_MODEL"
    TON_N = "#TOP_N_Wait_Events"
    TOP_SCHEMA_SQL = "#TOP_Schema_SQL"
    LITERAL_SQL = "#Literal_SQL"
    METRIC = "#METRIC"
    RAC = "#RAC"
    MEMORY = "#MEMORY"
    TOP_3 = "#TOP_3_Wait_Events"

    def __init__(self, logger):
        super().__init__(logger=logger)
        self.ot: OracleTarget = None
        self.sql_path = None
        self.position = None
        self.presentiton_path = 'sample_test.pptx'
        self.presentation = Presentation(self.presentiton_path)
        self.instance_number = None
        self.instance_name = None

    def main_process(self):

        self.logger.info("performance_analyzer.pptx")
        self._insert_extend_target_data()

        self.instance_number = self._extract_instance_info('instance_number')
        self.instance_name = self._extract_instance_info('instance_name')

        self.sql_path = f"{self.config['home']}/" + SystemConstants.CHART_SQL
        self.position = Config("position").get_config()

        # self._execute_time_model()
        # self._execute_db_system()
        # self._execute_metric()
        # self._execute_rac()
        # self._execute_top_n_wait_events()
        self._execute_top_3_wait_events()
        # self._execute_top_schema_sql()
        # self._execute_memory()
        # self._execute_literal_sql()
        # SlideManager.delete_slide(self.presentation)
        # self.presentation.save("instance1.pptx")

    def _execute_delete_slide(self):

        xml_slides = self.presentation.slides._sldIdLst
        xml_slides_list = list(xml_slides)
        for idx,slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.has_text_frame and re.search(r"^#",  shape.text_frame.text):
                    xml_slides.remove(xml_slides_list[idx])

    def _insert_extend_target_data(self):
        """
        DB 확장 타겟(실제 분석 대상 DB) 데이터 저장 함수
        """
        extend_target_repo_list = self.config["maxgauge_repo"].get("extend_target_repo", [])

        for extend_target_repo in extend_target_repo_list:
            extend_target_repo["analysis_target_type"] = self.config["maxgauge_repo"]["analysis_target_type"]  # oracle

            if str(self.config["maxgauge_repo"]["analysis_target_type"]).lower() == DbTypeConstants.ORACLE:
                if self.ot is None:
                    self.ot = OracleTarget(self.logger, self.config)

                self.ot.set_extend_target_config(extend_target_repo)
                self.ot.init_process()

    def _extract_instance_info(self, target):
        query = PptSql.SELECT_INSTANCE_NUMBER
        for df in self.ot.get_data_by_query(query):
            value_list = df[target].unique().tolist()
            value_list.sort(reverse=False)
            return value_list

    def _extract_metric_name_list(self):

        metric_name_list = [
            "Host CPU Utilization (%)",
            "Average Active Sessions",
            "Executions Per Sec",
            "User Transaction Per Sec",
            "Logical Reads Per Sec",
            "Physical Reads Per Sec",
            "Hard Parse Count Per Sec",
        ]

        config_report = Config("report").get_config()
        metric_name_list.extend(config_report["sys_metric"])

        return metric_name_list

    def _convert_sql_to_df(self, sql_path, filename,event_name="", inst_num=""):

        """
        sql query문을 dataframe 형태로 변환
        """

        metric_name_list = self._extract_metric_name_list()
        unpack_metric_name_list = str(metric_name_list)[1:-1]

        s_date, e_date = DateUtils.get_each_date_by_interval2(
            self.config["args"]["s_date"], self.config["args"]["interval"], arg_fmt="%Y-%m-%d"
        )
        date_dict = {"StartDate": s_date, "EndDate": e_date, "Metric_Name": unpack_metric_name_list,"EVENT_NAME":event_name, "INST_NUM":inst_num}

        query = SystemUtils.get_file_content_in_path(sql_path, filename + ".txt")
        date_query = SqlUtils.sql_replace_to_dict(query, date_dict)
        for df in self.ot.get_data_by_query(date_query):
            df.columns = [i.upper() for i in df.columns]
            df = df.fillna(0)
            return df

    def is_even(self, i):
        return i % 2 == 0

    def _execute_time_model(self):

        self.logger.info("time_model.pptx")

        left_tp = SlideManager.convert_inches_to_data(self.position['time_model']['left_table'])
        right_tp = SlideManager.convert_inches_to_data(self.position['time_model']['right_table'])
        column_width_inches = [Inches(i) for i in self.position['time_model']['column_width_inches']]

        time_model_df = self._convert_sql_to_df(self.sql_path, 'TIME_MODEL')
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TIME_MODEL)

        k = 1

        for i in range(len(self.instance_number)):

            if i == 0:

                SlideManager.create_text_box(target_slide, self.position['top_menu']['left_top_text_position'],
                                             "DB Workload 1",
                                             PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                             PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._time_model_detail(time_model_df, self.instance_number[i:i + 2], self.instance_name[i:i+2], target_slide,left_tp,right_tp,column_width_inches)

            if i > 1 and i % 2 == 0:

                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, k)
                SlideManager.create_text_box(copied_slide, self.position['top_menu']['left_top_text_position'],
                                             "DB Workload 1",
                                             PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                             PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._time_model_detail(time_model_df, self.instance_number[i:i + 2], self.instance_name[i:i+2], copied_slide,left_tp,right_tp,column_width_inches)
                k += 1

    def _time_model_detail(self, df, instance_num_range, instance_name_range, slide,left_table_position,right_table_position,column_width_inches):

        for i in range(len(instance_num_range)):

            inst_df = df[df['INSTANCE_NUMBER'] == instance_num_range[i]]
            inst_num_drop_df = inst_df.drop(columns=['INSTANCE_NUMBER'])

            rec_position = self.position['time_model']['left_rectangle_position'] if self.is_even(i) else self.position['time_model']['right_rectangle_position']
            instance_name_position = self.position['time_model']['left_position'] if self.is_even(i) else self.position['time_model']['right_position']
            tp = left_table_position if self.is_even(i) else right_table_position

            SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, rec_position,
                                      PerformanceAnalyzer.YELLOW_COLOR, PerformanceAnalyzer.YELLOW_COLOR, 0)

            SlideManager.create_text_box(slide, instance_name_position, instance_name_range[i],
                                         PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                         PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
            SlideManager.make_table(inst_num_drop_df, tp, slide,column_width_inches)

    def _execute_db_system(self):

        self.logger.info("db_system.pptx")

        column_width_inches = [Inches(i) for i in self.position['db_system']['column_width_inches']]
        tp = SlideManager.convert_inches_to_data(self.position['db_system']['table_position'])
        df = self._convert_sql_to_df(self.sql_path, "DB_SYSTEM")
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, "#DB_SYSTEM")

        k = 1

        for i in range(len(self.instance_number)):
            if i == 0:
                self._db_system_detail(df, self.instance_number[0:2], target_slide,column_width_inches,tp,i)

            elif i % 2 == 0:
                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, k)
                self._db_system_detail(df, self.instance_number[i: i+2], copied_slide,column_width_inches,tp,i)
                k += 1

    def _db_system_detail(self, df, instance_num_range, slide,column_width_inches,tp,i):

        list_df = []
        for inst_num in instance_num_range:
            inst_df = df[df['INSTANCE_NUMBER'] == inst_num]
            list_df.append(inst_df)
        preprocessed_df = self._arrange_db_system(list_df, i)

        SlideManager.create_text_box(slide, self.position['top_menu']['left_top_text_position'], "DB SYSTEM",
                                     PerformanceAnalyzer.LARGE_FONT_SIZE, "나눔스퀘어 네오 Light", True, 1.0)
        SlideManager.make_table(preprocessed_df, tp, slide, column_width_inches)

    def _arrange_db_system(self, list_df,i):

        dfs_to_merge = []
        new_row = pd.DataFrame({"CATEGORY": "Business"}, index=[0])

        for df in list_df:

            inst_num_drop_df = df.drop(columns=['INSTANCE_NUMBER'])
            concat_df = pd.concat([new_row, inst_num_drop_df]).reset_index(drop=True)

            arranged_df = concat_df.replace(['HOST_NAME','INSTANCE_NAME','PLATFORM_NAME','BANNER','NUM_CPUS',
                                      'PHYSICAL_MEMORY_BYTES','Buffer Cache Size','Shared Pool Size','Large Pool Size',
                                      'Java Pool Size','Streams Pool Size','pga_aggregate_limit'],
                                     ['Hostname','Instance Name','OS Version','GI & DB Version','OS CPU','OS Memory',
                                      'Buffer Cache','Shared Pool','Large Pool','Java Pool','Streams Pool','PGA'])

            arranged_df.loc[1:12,'CATEGORY'] = ['Hostname', 'OS Version', 'Instance Name', 'GI & DB Version','OS CPU', 'OS Memory',
                                  'Buffer Cache','Shared Pool','Java Pool','Large Pool','Streams Pool','PGA']

            arranged_df = arranged_df.fillna(0)
            arranged_df.rename(columns={'VALUE':f"NODE{i+1}"}, inplace=True)

            dfs_to_merge.append(arranged_df)

        if len(dfs_to_merge) == 1:
            return arranged_df

        else:
            merge_df = pd.merge(dfs_to_merge[0], dfs_to_merge[1], on='CATEGORY', how='outer')
            merge_df.rename(columns={f"NODE{i+1}_x": f"NODE{i+1}",f"NODE{i+1}_y": f"NODE{i+2}"}, inplace=True)
            return merge_df

    def _execute_metric(self):

        self.logger.info("metric.pptx")
        df = self._convert_sql_to_df(self.sql_path, "METRIC")
        metric_name_list = self._extract_metric_name_list()
        num_slide, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.METRIC)

        for i in range(len(metric_name_list)):

            if i == 0:
                self._metric_detail(target_slide, metric_name_list[i], df)

            else:
                _, copied_slide = SlideManager.add_slide(self.presentation, num_slide, i)
                self._metric_detail(copied_slide, metric_name_list[i], df)

    def _metric_detail(self, slide, metric_name, df):

        SlideManager.create_text_box(slide, self.position['top_menu']['left_top_text_position'], metric_name, PerformanceAnalyzer.LARGE_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  self.position['metric']['avg_rectangle_position'], PerformanceAnalyzer.YELLOW_COLOR,
                                  PerformanceAnalyzer.YELLOW_COLOR, 0)

        SlideManager.create_text_box(slide, self.position['metric']['avg_text_position'], "AVG",PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  self.position['metric']['max_rectangle_position'], PerformanceAnalyzer.YELLOW_COLOR,
                                  PerformanceAnalyzer.YELLOW_COLOR, 0)

        SlideManager.create_text_box(slide, self.position['metric']['max_text_position'], "MAX",PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        metric_df = df[df["METRIC_NAME"] == metric_name]
        copied_df = metric_df.copy()
        preprocessed_metric_df = self._set_df_date_time(copied_df)
        upper_chart_data, down_chart_data, upper_scale, down_scale = self._insert_chart(preprocessed_metric_df,'AG','MX')

        SlideManager.set_y_axis_max_value(slide,upper_chart_data, self.position['metric']['avg_coordinate'], down_scale)
        SlideManager.set_y_axis_max_value(slide,down_chart_data, self.position['metric']['max_coordinate'], down_scale)

        self._chart_style(slide)
        self._set_bottom_text(slide)

    def _execute_rac(self):

        self.logger.info("RAC.pptx")
        if len(self.instance_name) != 1:
            df = self._convert_sql_to_df(self.sql_path, "RAC")
            preprocessed_df = self._set_df_date_time(df)
            num_slide, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.RAC)

            SlideManager.create_text_box(target_slide,self.position['top_menu']['left_top_text_position'],
                                         "RAC Interconnect traffic", PerformanceAnalyzer.LARGE_FONT_SIZE,
                                         PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                         PerformanceAnalyzer.BASE_LINE_SPACE)

            SlideManager.create_shape(target_slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                      self.position['rac']['rec1_position'],
                                      PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)

            SlideManager.create_shape(target_slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                      self.position['rac']['rec2_position'],
                                      PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)

            SlideManager.create_text_box(target_slide, self.position['rac']['sent_position'],
                                         "DB Block Sent",PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                                         PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                         PerformanceAnalyzer.BASE_LINE_SPACE)

            SlideManager.create_text_box(target_slide, self.position['rac']['received_position'],
                                         "DB Block Received",PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                                         PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                         PerformanceAnalyzer.BASE_LINE_SPACE)

            upper_chart_data, down_chart_data, upper_scale, down_scale = self._insert_chart(preprocessed_df, 'BYTES_SENTPSEC', 'BYTES_RECEIVEDPSEC')
            SlideManager.set_y_axis_max_value(target_slide, upper_chart_data,self.position['rac']['sent_coordinate'], upper_scale)
            SlideManager.set_y_axis_max_value(target_slide, down_chart_data,self.position['rac']['received_coordinate'], down_scale)

            self._chart_style(target_slide)
            self._set_bottom_text(target_slide)

        else:
            pass

    def _extract_name_value_list(self):

        amm_memory_df = self._convert_sql_to_df(self.sql_path, "MEMORY_AMM")

        amm_memory_list = ['Physical memory','sga_target','db_block_size','pga_aggregate_target']
        name_value_list = []
        for name in amm_memory_list:
            filtered_row = amm_memory_df[amm_memory_df['NAME'] == name]
            amm_memory_value = filtered_row['VALUE'].iloc[0]

            if name == 'Physical memory':
                changed_name = 'Physical Memory'

            else:
                changed_name =name.upper()

            make_sentence = f'{changed_name} = {amm_memory_value}'
            name_value_list.append(make_sentence)

        return name_value_list

    def _execute_memory(self):

        self.logger.info("memory.pptx")

        left_tp = SlideManager.convert_inches_to_data(self.position['memory']['left_table'])
        right_tp = SlideManager.convert_inches_to_data(self.position['memory']['right_table'])
        column_width_inches = [Inches(i) for i in self.position['memory']['column_width_inches']]

        shared_memory_df = self._convert_sql_to_df(self.sql_path, "MEMORY_INFO")
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.MEMORY)

        k = 1

        for i in range(len(self.instance_number)):

            if i ==0:
                self._memory_detail(target_slide, shared_memory_df, self.instance_number[i:i+2], self.instance_name[i:i+2], left_tp, right_tp, column_width_inches)
                self._create_memory_blue_box(target_slide)

            if i >1 and i %2 ==0:
                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, k)
                self._memory_detail(copied_slide, shared_memory_df, self.instance_number[i:i+2], self.instance_name[i:i+2], left_tp, right_tp, column_width_inches)
                self._create_memory_blue_box(copied_slide)
                k+=1

    def _create_memory_blue_box(self, slide):


        blue_box = SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                             self.position['memory']['rectangle_position'],
                                             PerformanceAnalyzer.WHILE_COLOR,
                                             PerformanceAnalyzer.LIGHT_BLUE_COLOR, 0)

        blue_text_box = SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                                  self.position['memory']['rectangle2_position'],
                                                  PerformanceAnalyzer.WHILE_COLOR,
                                                  PerformanceAnalyzer.WHILE_COLOR, 0)

        SlideManager.create_text_frame(blue_text_box.text_frame.paragraphs[0],
                                       "오라클 AMM(Automatic Memory Management) 적용",
                                       "Roboto Light", PerformanceAnalyzer.MEDIUM_FONT_SIZE,
                                       PerformanceAnalyzer.LIGHT_BLUE_COLOR, PerformanceAnalyzer.BOLD_TRUE, 0)

        for content_text in self._extract_name_value_list():
            SlideManager.create_text_frame(blue_box.text_frame.add_paragraph(), content_text, "Roboto",
                                           PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.BLACK_COLOR,
                                           PerformanceAnalyzer.BOLD_FALSE, PerformanceAnalyzer.LONG_LINE_SPACE)

    def _memory_detail(self, slide, df, instance_num_range, instance_name_range, left_tp, right_tp, column_width_inches):

        for i in range(len(instance_num_range)):

            rec_position = self.position['memory']['ora1_rectangle_position'] if self.is_even(i) else self.position['memory']['ora2_rectangle_position']
            instance_name_position = self.position['memory']['oracle_1_position'] if self.is_even(i) else self.position['memory']['oracle_2_position']
            except_columns = ['INST_ID']

            inst_df = SlideManager.extract_specified_df(df, 'INST_ID', except_columns, instance_num_range[i])

            tp = left_tp if self.is_even(i) else right_tp

            SlideManager.create_text_box(slide, self.position['top_menu']['left_top_text_position'],
                                         "성능 분석 – 공유 메모리 점검", PerformanceAnalyzer.LARGE_FONT_SIZE,
                                         PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                         PerformanceAnalyzer.BASE_LINE_SPACE)

            SlideManager.create_text_box(slide, instance_name_position, instance_name_range[i],
                                         PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                         PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
            SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,rec_position,
                                      PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)
            SlideManager.make_table(inst_df, tp, slide, column_width_inches)

    def _execute_top_schema_sql(self):

        self.logger.info("top_schema_sql.pptx")
        top_schema_sql1 = self._convert_sql_to_df(self.sql_path, "TOP_Schema_SQL")
        top_schema_sql2 = self._convert_sql_to_df(self.sql_path, "TOP_Schema_SQL2")

        pd.set_option('display.max_columns', None)
        pd.set_option('display.max_rows', None)

        top_schema_sql2['SQL_TEXT'] = top_schema_sql2['SQL_TEXT'].str.replace(r'\s+', ' ').apply(lambda x: x.lstrip())
        top_schema_sql2['SQL_TEXT'] = top_schema_sql2['SQL_TEXT'].str.slice(0, 60)

        top_schema_sql2['SQL_TEXT'] = top_schema_sql2['SQL_TEXT'].astype(str)+"..."

        self._execute_top_schema_sql1(top_schema_sql1, top_schema_sql2)

    def _execute_top_schema_sql1(self, upper_df, down_df):

        upper_tp = SlideManager.convert_inches_to_data(self.position['top_schema_sql']['upper_table'])
        down_tp = SlideManager.convert_inches_to_data(self.position['top_schema_sql']['down_table'])
        upper_column_width_inches = [Inches(i) for i in self.position['top_schema_sql']['table1_column_inches']]
        down_column_width_inches = [Inches(i) for i in self.position['top_schema_sql']['table2_column_inches']]

        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TOP_SCHEMA_SQL)

        for i in range(len(self.instance_number)):

            if i == 0:
                SlideManager.create_text_box(target_slide, self.position['top_menu']['left_top_text_position'],
                                             "TOP Schema & SQL", PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                             PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._detail_top_schema_sql(target_slide, upper_df, down_df, upper_tp, down_tp, upper_column_width_inches, down_column_width_inches, i)

            if i >= 1:
                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, i)
                SlideManager.create_text_box(copied_slide, self.position['top_menu']['left_top_text_position'],
                                             "TOP Schema & SQL", PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                             PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._detail_top_schema_sql(copied_slide, upper_df, down_df, upper_tp, down_tp, upper_column_width_inches, down_column_width_inches, i)

    def _detail_top_schema_sql(self, slide, upper_df, down_df, upper_tp, down_tp, upper_column_width_inches, down_column_width_inches, i):


        up_except_col = ['INSTANCE_NUMBER','IOWAIT_TIME(SEC)', 'IOWAIT_RATIO(%)', 'RNK']
        down_except_col = ['INSTANCE_NUMBER','RNK']

        sql1_df = SlideManager.extract_specified_df(upper_df, 'INSTANCE_NUMBER', up_except_col, self.instance_number[i])
        sql2_df = SlideManager.extract_specified_df(down_df, 'INSTANCE_NUMBER', down_except_col, self.instance_number[i])
        SlideManager.make_table(sql1_df, upper_tp, slide, upper_column_width_inches)
        SlideManager.make_table(sql2_df, down_tp, slide, down_column_width_inches)

        SlideManager.create_text_box(slide,self.position['top_menu']['left_top_text_position'],"TOP Schema & SQL", PerformanceAnalyzer.LARGE_FONT_SIZE, "나눔스퀘어 네오 Light", True, 1.0)

        SlideManager.create_text_box(slide, self.position['top_schema_sql']['oracle_num_position'], self.instance_name[i],
                                     PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                     PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,self.position['top_schema_sql']['rectangle_position'],
                                  PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)

        self._set_bottom_text(slide)

    def _execute_top_n_wait_events(self):

        self.logger.info("top_n_wait_events.pptx")
        left_tp = SlideManager.convert_inches_to_data(self.position['top_n']['left_table'])
        right_tp = SlideManager.convert_inches_to_data(self.position['top_n']['right_table'])
        column_width_inches = [Inches(i) for i in self.position['top_n']['column_width_inches']]

        top_n_df = self._convert_sql_to_df(self.sql_path, 'TOP_N_Wait_Events')
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TON_N)

        k = 1

        for i in range(len(self.instance_number)):

            if i ==0:
                SlideManager.create_text_box(target_slide, self.position["top_menu"]["left_top_text_position"],"TOP-N Wait Events",
                                                 PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                                 PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._top_n_wait_events_detail(target_slide, top_n_df, self.instance_number[i:i+2], self.instance_name[i:i+2],left_tp, right_tp, column_width_inches )

            if i >1 and i%2 ==0:

                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, k)
                SlideManager.create_text_box(copied_slide, self.position["top_menu"]["left_top_text_position"],
                                             "TOP-N Wait Events",
                                             PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                             PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
                self._top_n_wait_events_detail(copied_slide, top_n_df, self.instance_number[i:i+2], self.instance_name[i:i+2], left_tp, right_tp, column_width_inches)
                k+=1

    def _top_n_wait_events_detail(self, slide, df, instance_num_range, instance_name_range, left_tp, right_tp, column_width_inches):

        for i in range(len(instance_num_range)):

            inst_df = df[df['INSTANCE_NUMBER'] == instance_num_range[i]]
            inst_num_drop_df = inst_df.drop(columns=['INSTANCE_NUMBER', 'RNK'])

            rec_position = self.position['top_n']['left_rectangle_position'] if self.is_even(i) else self.position['top_n']['right_rectangle_position']
            instance_name_position = self.position['top_n']['left_position'] if self.is_even(i) else self.position['top_n']['right_position']
            tp = left_tp if self.is_even(i) else right_tp

            SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, rec_position,
                                      PerformanceAnalyzer.YELLOW_COLOR, PerformanceAnalyzer.YELLOW_COLOR, 0)

            SlideManager.create_text_box(slide, instance_name_position, instance_name_range[i],
                                         PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                         PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
            SlideManager.make_table(inst_num_drop_df, tp, slide,column_width_inches)
            self._set_bottom_text(slide)

    def _execute_top_3_wait_events(self):

        self.logger.info("top_3_wait_events.pptx")
        df = self._convert_sql_to_df(self.sql_path, "TOP_N_Wait_Events")
        print(df)

        instance_dict = self._extract_top3_list(df)
        print(instance_dict)

        upper_tp = SlideManager.convert_inches_to_data(self.position['top_3']['upper_table'])
        down_tp = SlideManager.convert_inches_to_data(self.position['top_3']['down_table'])
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.TOP_3)

        k = 0
        for indx, (key, value) in enumerate(instance_dict.items()):

            inst_num = key.split("_")[-1]
            for event_name in value:
                sql1_df = self._extract_preprocessed_df(event_name, inst_num, "TOP_1_Wait_Events")
                sql2_df = self._extract_preprocessed_df(event_name, inst_num, "TOP_2_Wait_Events", ['INSTANCE_NUMBER'])
                sql3_df = self._extract_preprocessed_df(event_name, inst_num, "TOP_3_Wait_Events", ['INSTANCE_NUMBER'])
                if k == 0:
                    self._top3_detail(target_slide, event_name, self.instance_number[indx], self.instance_name[indx],
                                      sql1_df, sql2_df, sql3_df, upper_tp, down_tp)
                    k += 1

                elif k >= 1:
                    copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, k)
                    self._top3_detail(copied_slide, event_name, self.instance_number[indx],self.instance_name[indx],
                                      sql1_df, sql2_df, sql3_df, upper_tp, down_tp)
                    k+=1


    def _top3_detail(self, slide, event_name, inst_num, inst_name, sql1_df, sql2_df, sql3_df, upper_tp, down_tp):

        SlideManager.create_text_box(slide, self.position['top_menu']['left_top_text_position'],
                                     event_name, PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                     PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  self.position['top_3']['ora_rectangle_position'],
                                  PerformanceAnalyzer.YELLOW_COLOR, PerformanceAnalyzer.YELLOW_COLOR, 0)

        SlideManager.create_text_box(slide, self.position['top_3']['ora_position'],
                                     inst_name, PerformanceAnalyzer.DEFAULT_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        chart_data, chart_scale = SlideManager.insert_chart(sql1_df, 'DATE_TIME', 'VALPSEC', inst_name)
        SlideManager.set_y_axis_max_value(slide, chart_data, self.position['top_3']['chart_position'], chart_scale)
        self._top3_chart_detail(slide, inst_num)
        SlideManager.make_table(sql2_df, upper_tp, slide)
        SlideManager.make_table(sql3_df, down_tp, slide)
        self._set_bottom_text(slide)





        #
        # for i in range(len(event_name_list)-1):
        #     index, copied_slide = SlideManager.add_slide(self.presentation, num_slide, i+1) #슬라이드 생성
        #     slide_list.append(copied_slide)
        #
        # slide_idx = 0
        #
        # for idx, (key, value) in enumerate(instance_dict.items()):
        #     instance_num = key.split("_")[-1]
        #     for event_name in value:
        #
        #         sql1_df = self._extract_preprocessed_df(event_name, instance_num, "TOP_1_Wait_Events")
        #         sql2_df = self._extract_preprocessed_df(event_name, instance_num, "TOP_2_Wait_Events", ['INSTANCE_NUMBER'])
        #         sql3_df = self._extract_preprocessed_df(event_name, instance_num, "TOP_3_Wait_Events", ['INSTANCE_NUMBER'])
        #
        #         SlideManager.create_text_box(slide_list[slide_idx],self.position['top_menu']['left_top_text_position'],
        #                                      event_name, PerformanceAnalyzer.LARGE_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
        #                                      PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
        #
        #         SlideManager.create_shape(slide_list[slide_idx], MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        #                                   self.position['top_3']['ora_rectangle_position'],
        #                                   PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)
        #
        #         SlideManager.create_text_box(slide_list[slide_idx], self.position['top_3']['ora_position'],
        #                                      self.instance_name[idx],PerformanceAnalyzer.DEFAULT_FONT_SIZE,
        #                                      PerformanceAnalyzer.FONT_NAME,PerformanceAnalyzer.BOLD_TRUE,
        #                                      PerformanceAnalyzer.BASE_LINE_SPACE)
        #
        #         chart_data, chart_scale = SlideManager.insert_chart(sql1_df, 'DATE_TIME', 'VALPSEC', self.instance_name[idx])
        #         SlideManager.set_y_axis_max_value(slide_list[slide_idx], chart_data, self.position['top_3']['chart_position'], chart_scale)
        #         self._top3_chart_detail(slide_list[slide_idx], instance_num)
        #         SlideManager.make_table(sql2_df, upper_table_position, slide_list[slide_idx])
        #         SlideManager.make_table(sql3_df, down_table_position, slide_list[slide_idx])
        #         self._set_bottom_text(slide_list[slide_idx])
        #         slide_idx += 1

    def _execute_literal_sql(self):

        self.logger.info("literalsql.pptx")
        slide_num, target_slide = SlideManager.read_slide(self.presentation.slides, PerformanceAnalyzer.LITERAL_SQL)
        column_width_inches = [Inches(i) for i in self.position['literal_sql']['column_width_inches']]
        tp = SlideManager.convert_inches_to_data(self.position['literal_sql']['table_position'])

        for i in range(len(self.instance_number)):

            if i == 0:
                self._literal_sql_detail(target_slide, tp, column_width_inches, self.instance_number[i], self.instance_name[i])

            if i >= 1:
                copied_num, copied_slide = SlideManager.add_slide(self.presentation, slide_num, i)
                self._literal_sql_detail(copied_slide, tp, column_width_inches,self.instance_number[i], self.instance_name[i])

    def _literal_sql_detail(self, slide, tp, column_width_inches, inst_num, instance_name):

        df = self._convert_sql_to_df(self.sql_path, filename="Literal_SQL", inst_num=str(inst_num))

        SlideManager.create_text_box(slide, self.position["top_menu"]["left_top_text_position"],
                                     "Top Literal SQL (TOP 20 Info)", PerformanceAnalyzer.LARGE_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.create_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  self.position['literal_sql']['ora_rectangle_position'],
                                  PerformanceAnalyzer.YELLOW_COLOR,PerformanceAnalyzer.YELLOW_COLOR, 0)

        SlideManager.create_text_box(slide, self.position['literal_sql']['ora_position'],
                                     instance_name, PerformanceAnalyzer.DEFAULT_FONT_SIZE,
                                     PerformanceAnalyzer.FONT_NAME, PerformanceAnalyzer.BOLD_TRUE,
                                     PerformanceAnalyzer.BASE_LINE_SPACE)

        SlideManager.make_table(df, tp, slide, column_width_inches)
        self._set_bottom_text(slide)

    def _top3_chart_detail(self, slide, instance_num):

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                SlideManager.set_chart_style(chart)
                plot = chart.plots[0]

                for index, series in enumerate(plot.series):
                    for point,label in enumerate(series.points):
                        if int(instance_num) == 1:
                            SlideManager.set_label_style(label,PerformanceAnalyzer.LIGHT_RED_COLOR, Pt(8), PerformanceAnalyzer.FONT_NAME,
                                                         True, XL_DATA_LABEL_POSITION.ABOVE)
                            SlideManager.set_chart_marker_style(series, XL_MARKER_STYLE.CIRCLE, PerformanceAnalyzer.LIGHT_RED_COLOR,
                                                                PerformanceAnalyzer.WHILE_COLOR)
                            SlideManager.set_chart_line_style(series, PerformanceAnalyzer.LIGHT_RED_COLOR, Pt(2.25))

                        if int(instance_num) == 2:
                            SlideManager.set_label_style(label, PerformanceAnalyzer.BLUE_GRAY_COLOR, Pt(8), PerformanceAnalyzer.FONT_NAME,
                                                         True, XL_DATA_LABEL_POSITION.ABOVE)
                            SlideManager.set_chart_marker_style(series, XL_MARKER_STYLE.CIRCLE, PerformanceAnalyzer.BLUE_GRAY_COLOR,
                                                                PerformanceAnalyzer.WHILE_COLOR)
                            SlideManager.set_chart_line_style(series, PerformanceAnalyzer.BLUE_GRAY_COLOR, Pt(2.25))

    def _extract_top3_list(self, df):

        instance_dict = {}
        for instance_num in df["INSTANCE_NUMBER"].unique():
            instance_num_df = df[df["INSTANCE_NUMBER"] == instance_num]
            top_3_df = instance_num_df.iloc[1:4]
            top3_3_name = list(top_3_df['EVENT_NAME'])
            instance_num_name = f"instance_{instance_num}"
            instance_dict[instance_num_name] = top3_3_name

        return instance_dict

    def _extract_preprocessed_df(self,event_name, instance_num, sql_filename, except_col=[]):

        df = self._convert_sql_to_df(self.sql_path, sql_filename, event_name)
        preprocessed_df = self._set_df_date_time(df)
        event_name_df = preprocessed_df[preprocessed_df['INSTANCE_NUMBER'] == int(instance_num)]
        event_name_df = event_name_df.drop(columns=except_col)

        if event_name_df.empty:
            event_name_df = pd.DataFrame(0, index=range(3), columns=event_name_df.columns)
        return event_name_df

    def _set_bottom_text(self, copied_slide):

        bottom_text = "RAC 의 경우 Fail CPU 사용률은 1,2번 모두 30% 이하로 안정적이며 특정 이벤트 및 I/O 급증으로 인한 CPU 증가 현상은 없습니다.\n" \
                      "Over 를 대비하여 각 인스턴스의 CPU 사용률은 50% 이하로 유지되도록 권장하며 응답소 DB의 경우 적정 수준을 유지하고 있습니다."

        SlideManager.create_text_box(copied_slide, self.position['bottom_menu']['bottom_findings_position'], 'Findings',
                                     PerformanceAnalyzer.DEFAULT_FONT_SIZE,PerformanceAnalyzer.FONT_NAME,
                                     PerformanceAnalyzer.BOLD_TRUE, PerformanceAnalyzer.BASE_LINE_SPACE)
        SlideManager.create_text_box(copied_slide, self.position['bottom_menu']['bottom_text_position'], bottom_text,
                                     PerformanceAnalyzer.DEFAULT_FONT_SIZE, PerformanceAnalyzer.FONT_NAME,
                                     PerformanceAnalyzer.BOLD_FALSE,PerformanceAnalyzer.LONG_LINE_SPACE)
        SlideManager.create_connector(copied_slide, self.position['bottom_menu']['straight_position'])
        SlideManager.create_shape(copied_slide, MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
                                  self.position['bottom_menu']['right_triangle_position'],
                                  PerformanceAnalyzer.DARK_GRAY_COLOR, PerformanceAnalyzer.DARK_GRAY_COLOR, -90)

    def _chart_style(self, slide):

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                SlideManager.set_chart_style(chart)
                plot = chart.plots[0]
                for index, series in enumerate(plot.series):
                    for point, label in enumerate(series.points):
                        if index == 0:
                            self._chart_style_detail(series, label, PerformanceAnalyzer.LIGHT_RED_COLOR)

                        elif index == 1:
                            self._chart_style_detail(series, label, PerformanceAnalyzer.BLUE_GRAY_COLOR)

                        else:
                            self._chart_style_detail(series, label, PerformanceAnalyzer.GREEN_COLOR)

    def _chart_style_detail(self, series, label, color):
        SlideManager.set_label_style(label, color, Pt(8), PerformanceAnalyzer.FONT_NAME, True,
                                     XL_DATA_LABEL_POSITION.ABOVE)
        SlideManager.set_chart_line_style(series, color, Pt(2.25))
        SlideManager.set_chart_marker_style(series, XL_MARKER_STYLE.CIRCLE, color,
                                            PerformanceAnalyzer.WHILE_COLOR)
    def _insert_chart(self, df,col1,col2):

        upper_chart_data = CategoryChartData()
        down_chart_data = CategoryChartData()

        upper_max_list = []
        down_max_list = []

        for idx, instance_num in enumerate(df['INSTANCE_NUMBER'].unique()):
            extract_df = df[df['INSTANCE_NUMBER'] == instance_num]
            chart_categories = extract_df['DATE_TIME'].tolist()

            upper_max_score = extract_df[col1].max()
            down_max_score = extract_df[col2].max()

            upper_col_tuple = tuple(extract_df[col1].tolist())
            down_col_tuple = tuple(extract_df[col2].tolist())

            upper_max_list.append(upper_max_score)
            down_max_list.append(down_max_score)

            upper_chart_data.categories = chart_categories
            upper_chart_data.add_series(f'ORCLDB_{instance_num}',upper_col_tuple)

            down_chart_data.categories = chart_categories
            down_chart_data.add_series(f'ORCLDB_{instance_num}', down_col_tuple)

        upper_scale = SlideManager.make_max_value(upper_max_list)
        down_scale = SlideManager.make_max_value(down_max_list)

        return upper_chart_data, down_chart_data, upper_scale, down_scale

    def _set_df_date_time(self, df):

        if 'DATE_TIME' in df.columns:
            df["DATE_TIME"] = df["DATE_TIME"].astype(str) + ":00"
            df["DATE_TIME"] = df["DATE_TIME"].apply(lambda x: "\n".join(x.split()))

        return df










